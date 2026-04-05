"""AI Presentations — generate slide decks via LLM, export to PDF/PPTX."""

import io
import json
import logging
import re
import uuid

import httpx
from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import Response, StreamingResponse
from pydantic import BaseModel, Field
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.config import get_settings
from app.database import get_db
from app.middleware.auth import get_current_user
from app.models.generation import Generation
from app.models.user import User
from app.services.ai_router import get_openrouter_model, DEFAULT_MODEL

logger = logging.getLogger(__name__)
router = APIRouter(prefix="/api/presentations", tags=["presentations"])

COST_USD = 0.421  # ~40 RUB / 95


# ─── Schemas ───

class GenerateRequest(BaseModel):
    topic: str
    slides_count: int = Field(default=10, ge=3, le=30)
    style: str = "modern"
    audience: str | None = None
    detail_level: str = Field(default="medium", pattern="^(short|medium|detailed)$")
    language: str = "ru"
    model_id: str | None = None
    project_id: str | None = None


class ExportPDFRequest(BaseModel):
    generation_id: str


# ─── Themes ───

THEMES = {
    "modern": {
        "bg": "#ffffff", "text": "#1a1a2e", "accent": "#6c63ff",
        "heading": "#1a1a2e", "bullet": "#6c63ff",
        "slide_bg_title": "linear-gradient(135deg, #6c63ff 0%, #3f3d9e 100%)",
        "title_text": "#ffffff",
    },
    "minimal": {
        "bg": "#fafafa", "text": "#333333", "accent": "#000000",
        "heading": "#111111", "bullet": "#888888",
        "slide_bg_title": "#111111",
        "title_text": "#ffffff",
    },
    "corporate": {
        "bg": "#ffffff", "text": "#2c3e50", "accent": "#2980b9",
        "heading": "#1a252f", "bullet": "#2980b9",
        "slide_bg_title": "linear-gradient(135deg, #2980b9 0%, #1a5276 100%)",
        "title_text": "#ffffff",
    },
    "creative": {
        "bg": "#fff8f0", "text": "#2d2d2d", "accent": "#ff6b6b",
        "heading": "#2d2d2d", "bullet": "#ff6b6b",
        "slide_bg_title": "linear-gradient(135deg, #ff6b6b 0%, #ffa07a 100%)",
        "title_text": "#ffffff",
    },
    "bold": {
        "bg": "#0f0f0f", "text": "#e0e0e0", "accent": "#ffcc00",
        "heading": "#ffffff", "bullet": "#ffcc00",
        "slide_bg_title": "linear-gradient(135deg, #ffcc00 0%, #ff9900 100%)",
        "title_text": "#0f0f0f",
    },
    "dark": {
        "bg": "#1e1e2e", "text": "#cdd6f4", "accent": "#89b4fa",
        "heading": "#cdd6f4", "bullet": "#89b4fa",
        "slide_bg_title": "linear-gradient(135deg, #313244 0%, #1e1e2e 100%)",
        "title_text": "#cdd6f4",
    },
}


# ─── Helpers ───

def _parse_slides_json(text: str) -> list[dict]:
    """Extract JSON array from LLM response, handling markdown code blocks."""
    # Strip markdown code fences
    cleaned = re.sub(r"^```(?:json)?\s*", "", text.strip())
    cleaned = re.sub(r"\s*```$", "", cleaned.strip())

    # Try to find JSON array
    match = re.search(r"\[[\s\S]*\]", cleaned)
    if not match:
        raise ValueError("No JSON array found in response")

    return json.loads(match.group(0))


def _build_prompt(req: GenerateRequest) -> str:
    detail_map = {
        "short": "2-3 bullet points per slide, very concise",
        "medium": "4-5 bullet points per slide, moderate detail",
        "detailed": "5-7 bullet points per slide, comprehensive",
    }
    audience_hint = f"\nTarget audience: {req.audience}" if req.audience else ""

    return f"""Create a presentation on the topic: "{req.topic}"

Number of slides: {req.slides_count}
Detail level: {detail_map[req.detail_level]}
Language: {req.language}{audience_hint}

Return ONLY a valid JSON array (no markdown, no explanation) of slides in this format:
[
  {{
    "title": "Slide title",
    "bullets": ["Point 1", "Point 2", "Point 3"],
    "notes": "Speaker notes for this slide",
    "layout": "title"
  }},
  {{
    "title": "Content slide",
    "bullets": ["Point 1", "Point 2"],
    "notes": "Notes...",
    "layout": "content"
  }}
]

Available layouts: "title" (for the first/section slides), "content" (regular), "section-header" (divider), "quote" (for quotes), "two-column" (split content).

The first slide MUST use "title" layout.
For "quote" layout, put the quote as the first bullet and attribution as the second.
For "two-column" layout, put left-column bullets first, then add "---" as separator, then right-column bullets.

Make the content insightful, well-structured, and engaging. Each slide should have a clear purpose."""


async def _find_user(db: AsyncSession, tg_user: dict) -> User:
    """Locate the User row by db_id or telegram_id."""
    db_id = tg_user.get("db_id")
    if db_id:
        result = await db.execute(select(User).where(User.id == db_id).with_for_update())
    else:
        result = await db.execute(select(User).where(User.telegram_id == tg_user["id"]).with_for_update())
    user = result.scalar_one_or_none()
    if not user:
        raise HTTPException(404, "Пользователь не найден")
    return user


# ─── HTML rendering ───

def render_presentation_html(slides: list[dict], style: str, title: str) -> str:
    """Render slides as printable HTML (960x540 per slide, inline styles)."""
    theme = THEMES.get(style, THEMES["modern"])
    html_slides = []

    for idx, slide in enumerate(slides, 1):
        layout = slide.get("layout", "content")
        slide_title = slide.get("title", "")
        bullets = slide.get("bullets", [])
        image_url = slide.get("image_url", "") or ""

        # Build slide body based on layout
        if layout == "title":
            bg = theme["slide_bg_title"]
            bg_css = f"background: {bg};" if bg.startswith("linear") else f"background-color: {bg};"
            body = f"""
            <div style="{bg_css} width:960px; height:540px; display:flex;
                         flex-direction:column; justify-content:center; align-items:center;
                         text-align:center; padding:60px; box-sizing:border-box;
                         page-break-after:always; position:relative;">
                <h1 style="color:{theme['title_text']}; font-size:36px; margin:0 0 20px 0;
                           font-family:'Segoe UI',Arial,sans-serif; font-weight:700;">{_esc(slide_title)}</h1>
                {''.join(f'<p style="color:{theme["title_text"]}; opacity:0.85; font-size:20px; margin:4px 0; font-family:Segoe UI,Arial,sans-serif;">{_esc(b)}</p>' for b in bullets)}
                <div style="position:absolute; bottom:16px; right:24px; color:{theme['title_text']}; opacity:0.5; font-size:12px;">{idx}</div>
            </div>"""

        elif layout == "section-header":
            bg = theme["slide_bg_title"]
            bg_css = f"background: {bg};" if bg.startswith("linear") else f"background-color: {bg};"
            body = f"""
            <div style="{bg_css} width:960px; height:540px; display:flex;
                         flex-direction:column; justify-content:center; align-items:center;
                         text-align:center; padding:60px; box-sizing:border-box;
                         page-break-after:always; position:relative;">
                <h2 style="color:{theme['title_text']}; font-size:32px; margin:0;
                           font-family:'Segoe UI',Arial,sans-serif; font-weight:700;">{_esc(slide_title)}</h2>
                <div style="position:absolute; bottom:16px; right:24px; color:{theme['title_text']}; opacity:0.5; font-size:12px;">{idx}</div>
            </div>"""

        elif layout == "quote":
            quote_text = bullets[0] if bullets else ""
            attribution = bullets[1] if len(bullets) > 1 else ""
            body = f"""
            <div style="background-color:{theme['bg']}; width:960px; height:540px; display:flex;
                        flex-direction:column; justify-content:center; align-items:center;
                        padding:80px; box-sizing:border-box; page-break-after:always; position:relative;">
                <div style="border-left:4px solid {theme['accent']}; padding-left:24px; max-width:700px;">
                    <p style="color:{theme['text']}; font-size:24px; font-style:italic;
                              font-family:'Segoe UI',Arial,sans-serif; line-height:1.5; margin:0 0 16px 0;">&ldquo;{_esc(quote_text)}&rdquo;</p>
                    <p style="color:{theme['accent']}; font-size:18px; font-family:'Segoe UI',Arial,sans-serif; margin:0;">— {_esc(attribution)}</p>
                </div>
                <div style="position:absolute; bottom:16px; right:24px; color:{theme['text']}; opacity:0.4; font-size:12px;">{idx}</div>
            </div>"""

        elif layout == "two-column":
            sep_idx = None
            for i, b in enumerate(bullets):
                if b.strip() == "---":
                    sep_idx = i
                    break
            left = bullets[:sep_idx] if sep_idx is not None else bullets
            right = bullets[sep_idx + 1:] if sep_idx is not None else []
            left_html = "".join(
                f'<li style="color:{theme["text"]}; font-size:18px; margin:6px 0; line-height:1.4;">{_esc(b)}</li>' for b in left
            )
            right_html = "".join(
                f'<li style="color:{theme["text"]}; font-size:18px; margin:6px 0; line-height:1.4;">{_esc(b)}</li>' for b in right
            )
            if image_url:
                all_bullets_html = "".join(
                    f'<li style="color:{theme["text"]}; font-size:18px; margin:6px 0; line-height:1.4;">{_esc(b)}</li>'
                    for b in bullets if b.strip() != "---"
                )
                body = f"""
            <div style="background-color:{theme['bg']}; width:960px; height:540px; display:flex;
                        box-sizing:border-box; page-break-after:always; position:relative;">
                <div style="width:60%; padding:48px 40px 48px 56px;">
                    <h2 style="color:{theme['heading']}; font-size:28px; margin:0 0 24px 0;
                               font-family:'Segoe UI',Arial,sans-serif; font-weight:700;">{_esc(slide_title)}</h2>
                    <ul style="list-style:none; padding:0; margin:0;">
                        {all_bullets_html}
                    </ul>
                </div>
                <div style="width:40%; display:flex; align-items:center; justify-content:center; padding:24px;">
                    <img src="{image_url}" style="max-width:100%; max-height:90%; object-fit:contain; border-radius:8px;" />
                </div>
                <div style="position:absolute; bottom:16px; right:24px; color:{theme['text']}; opacity:0.4; font-size:12px;">{idx}</div>
            </div>"""
            else:
                body = f"""
            <div style="background-color:{theme['bg']}; width:960px; height:540px; padding:48px 56px;
                        box-sizing:border-box; page-break-after:always; position:relative;">
                <h2 style="color:{theme['heading']}; font-size:28px; margin:0 0 24px 0;
                           font-family:'Segoe UI',Arial,sans-serif; font-weight:700;">{_esc(slide_title)}</h2>
                <div style="display:flex; gap:40px;">
                    <ul style="flex:1; list-style:none; padding:0; margin:0;">
                        <li style="color:{theme['accent']}; font-weight:700; font-size:14px; margin-bottom:8px; text-transform:uppercase;">Column 1</li>
                        {left_html}
                    </ul>
                    <ul style="flex:1; list-style:none; padding:0; margin:0;">
                        <li style="color:{theme['accent']}; font-weight:700; font-size:14px; margin-bottom:8px; text-transform:uppercase;">Column 2</li>
                        {right_html}
                    </ul>
                </div>
                <div style="position:absolute; bottom:16px; right:24px; color:{theme['text']}; opacity:0.4; font-size:12px;">{idx}</div>
            </div>"""

        else:  # content (default)
            bullets_html = "".join(
                f'<li style="color:{theme["text"]}; font-size:18px; margin:8px 0; line-height:1.5;'
                f' list-style-type:none; padding-left:16px; position:relative;">'
                f'<span style="color:{theme["bullet"]}; position:absolute; left:0;">●</span> {_esc(b)}</li>'
                for b in bullets
            )
            if image_url:
                body = f"""
            <div style="background-color:{theme['bg']}; width:960px; height:540px; display:flex;
                        box-sizing:border-box; page-break-after:always; position:relative;">
                <div style="width:60%; padding:48px 40px 48px 56px;">
                    <h2 style="color:{theme['heading']}; font-size:28px; margin:0 0 24px 0;
                               font-family:'Segoe UI',Arial,sans-serif; font-weight:700;">{_esc(slide_title)}</h2>
                    <ul style="padding:0; margin:0;">{bullets_html}</ul>
                </div>
                <div style="width:40%; display:flex; align-items:center; justify-content:center; padding:24px;">
                    <img src="{image_url}" style="max-width:100%; max-height:90%; object-fit:contain; border-radius:8px;" />
                </div>
                <div style="position:absolute; bottom:16px; right:24px; color:{theme['text']}; opacity:0.4; font-size:12px;">{idx}</div>
            </div>"""
            else:
                body = f"""
            <div style="background-color:{theme['bg']}; width:960px; height:540px; padding:48px 56px;
                        box-sizing:border-box; page-break-after:always; position:relative;">
                <h2 style="color:{theme['heading']}; font-size:28px; margin:0 0 24px 0;
                           font-family:'Segoe UI',Arial,sans-serif; font-weight:700;">{_esc(slide_title)}</h2>
                <ul style="padding:0; margin:0;">{bullets_html}</ul>
                <div style="position:absolute; bottom:16px; right:24px; color:{theme['text']}; opacity:0.4; font-size:12px;">{idx}</div>
            </div>"""

        html_slides.append(body)

    return f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"><title>{_esc(title)}</title>
<style>
@page {{ size: 960px 540px; margin: 0; }}
body {{ margin: 0; padding: 0; }}
</style></head><body>
{''.join(html_slides)}
</body></html>"""


def _esc(text: str) -> str:
    """Escape HTML special characters."""
    return (
        str(text)
        .replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )


# ─── Endpoints ───

@router.post("/generate")
async def generate_presentation(
    req: GenerateRequest,
    tg_user: dict = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Generate a slide deck via LLM."""
    settings = get_settings()
    tg_id = tg_user["id"]

    # Find user and check balance
    user = await _find_user(db, tg_user)
    balance = float(user.balance_usd or 0)
    tier = user.subscription_tier or "free"
    has_subscription = tier in ("mini", "max", "max-pro")

    # Expensive models cost 3x
    EXPENSIVE_MODELS = {"gpt-4.1", "claude-sonnet-4", "claude-sonnet-4.5", "deepseek-r1"}
    model_id = req.model_id or DEFAULT_MODEL
    multiplier = 3 if model_id in EXPENSIVE_MODELS else 1
    actual_cost = COST_USD * multiplier

    if not has_subscription and balance < actual_cost:
        raise HTTPException(402, {
            "error": "insufficient_balance",
            "message": f"Недостаточно средств для генерации презентации (≈{round(actual_cost * 95):.0f}₽)",
            "required_usd": actual_cost,
            "balance_usd": balance,
        })

    # Deduct balance upfront (refund on failure)
    new_balance = balance
    if not has_subscription:
        user.balance_usd = round(balance - actual_cost, 6)
        new_balance = float(user.balance_usd)
        await db.flush()
    openrouter_model_id = get_openrouter_model(model_id)

    # Build prompt and call OpenRouter (non-streaming)
    prompt_text = _build_prompt(req)
    messages = [{"role": "user", "content": prompt_text}]

    try:
        async with httpx.AsyncClient(timeout=120.0) as client:
            resp = await client.post(
                "https://openrouter.ai/api/v1/chat/completions",
                headers={
                    "Authorization": f"Bearer {settings.openrouter_api_key}",
                    "HTTP-Referer": settings.webapp_url,
                    "X-Title": "Stone AI",
                    "Content-Type": "application/json",
                },
                json={
                    "model": openrouter_model_id,
                    "messages": messages,
                    "max_tokens": 8192,
                },
            )
            if resp.status_code != 200:
                raise RuntimeError(f"OpenRouter API error {resp.status_code}: {resp.text[:300]}")

            data = resp.json()
    except Exception as exc:
        # Refund on failure
        if not has_subscription:
            user.balance_usd = round(new_balance + actual_cost, 6)
            await db.flush()
        logger.error(f"Presentation generation failed: {exc}")
        raise HTTPException(502, f"Ошибка генерации: {exc}")

    # Extract content
    raw_content = data.get("choices", [{}])[0].get("message", {}).get("content", "")
    if not raw_content:
        if not has_subscription:
            user.balance_usd = round(new_balance + actual_cost, 6)
            await db.flush()
        raise HTTPException(502, "Модель вернула пустой ответ")

    # Parse slides JSON
    try:
        slides = _parse_slides_json(raw_content)
    except (json.JSONDecodeError, ValueError) as exc:
        if not has_subscription:
            user.balance_usd = round(new_balance + actual_cost, 6)
            await db.flush()
        logger.error(f"Failed to parse slides JSON: {exc}\nRaw: {raw_content[:500]}")
        raise HTTPException(502, "Не удалось разобрать ответ модели как JSON-презентацию")

    # Save generation
    gen = Generation(
        user_tg_id=tg_id,
        project_id=req.project_id,
        type="presentation",
        model=model_id,
        prompt=req.topic,
        result_text=json.dumps(slides, ensure_ascii=False),
        metadata_={"style": req.style, "slides_count": req.slides_count, "detail_level": req.detail_level, "language": req.language, "audience": req.audience},
        cost=actual_cost,
    )
    db.add(gen)
    await db.commit()
    await db.refresh(gen)

    cost_rub = round(actual_cost * 95, 2)

    return {
        "id": gen.id,
        "topic": req.topic,
        "slides": slides,
        "style": req.style,
        "generation_id": gen.id,
        "cost_rub": cost_rub,
        "balance_usd": new_balance,
    }


@router.get("/{generation_id}")
async def get_presentation(
    generation_id: str,
    tg_user: dict = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Retrieve a saved presentation by generation ID."""
    tg_id = tg_user["id"]

    result = await db.execute(
        select(Generation).where(
            Generation.id == generation_id,
            Generation.type == "presentation",
            Generation.user_tg_id == tg_id,
        )
    )
    gen = result.scalar_one_or_none()
    if not gen:
        raise HTTPException(404, "Презентация не найдена")

    slides = json.loads(gen.result_text) if gen.result_text else []
    meta = gen.metadata_ or {}

    return {
        "id": gen.id,
        "topic": gen.prompt,
        "slides": slides,
        "style": meta.get("style", "modern"),
        "model": gen.model,
        "generation_id": gen.id,
        "cost_rub": round((gen.cost or 0) * 95, 2),
        "created_at": gen.created_at.isoformat() if gen.created_at else None,
    }


@router.post("/export/html")
async def export_html(
    req: ExportPDFRequest,
    tg_user: dict = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Export presentation as printable HTML (client uses window.print → PDF)."""
    tg_id = tg_user["id"]

    result = await db.execute(
        select(Generation).where(
            Generation.id == req.generation_id,
            Generation.type == "presentation",
            Generation.user_tg_id == tg_id,
        )
    )
    gen = result.scalar_one_or_none()
    if not gen:
        raise HTTPException(404, "Презентация не найдена")

    slides = json.loads(gen.result_text) if gen.result_text else []
    meta = gen.metadata_ or {}
    style = meta.get("style", "modern")

    html = render_presentation_html(slides, style, gen.prompt or "Presentation")

    return Response(content=html, media_type="text/html")


# ─── PPTX Export ───

def _hex_to_rgb(hex_color: str):
    """Convert hex color to pptx RGBColor."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    h = hex_color.lstrip("#")
    if len(h) != 6:
        return RGBColor(0, 0, 0)
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _extract_gradient_color(val: str) -> str:
    """Extract first hex color from gradient or plain value."""
    match = re.search(r"#([0-9a-fA-F]{6})", val)
    return f"#{match.group(1)}" if match else val


def build_pptx(slides: list[dict], style: str, title: str) -> bytes:
    """Build a PPTX file from slide data."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    theme = THEMES.get(style, THEMES["modern"])
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    accent_color = _hex_to_rgb(theme["accent"])
    text_color = _hex_to_rgb(theme["text"])
    heading_color = _hex_to_rgb(theme["heading"])
    title_bg = _extract_gradient_color(theme.get("slide_bg_title", theme["accent"]))
    title_text_color = _hex_to_rgb(theme.get("title_text", "#ffffff"))
    bg_color = _hex_to_rgb(theme["bg"])
    bullet_color = _hex_to_rgb(theme.get("bullet", theme["accent"]))

    for slide_data in slides:
        layout = slide_data.get("layout", "content")
        slide_title = slide_data.get("title", "")
        bullets = slide_data.get("bullets", [])
        quote = slide_data.get("quote", "")
        author = slide_data.get("author", "")

        blank_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(blank_layout)

        # Background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        if layout in ("title", "section-header"):
            fill.fore_color.rgb = _hex_to_rgb(title_bg)
        else:
            fill.fore_color.rgb = bg_color

        if layout == "title":
            # Title slide
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(3))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_title
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = title_text_color
            p.alignment = PP_ALIGN.CENTER
            # Subtitle from first bullet
            if bullets:
                p2 = tf.add_paragraph()
                p2.text = bullets[0]
                p2.font.size = Pt(20)
                p2.font.color.rgb = RGBColor(255, 255, 255) if theme.get("title_text") == "#ffffff" else text_color
                p2.alignment = PP_ALIGN.CENTER
                p2.space_before = Pt(16)

        elif layout == "section-header":
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_title
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = title_text_color
            p.alignment = PP_ALIGN.CENTER

        elif layout == "quote":
            # Quote slide
            txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.333), Inches(4))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f'"{quote or slide_title}"'
            p.font.size = Pt(28)
            p.font.italic = True
            p.font.color.rgb = accent_color
            p.alignment = PP_ALIGN.CENTER
            if author:
                p2 = tf.add_paragraph()
                p2.text = f"— {author}"
                p2.font.size = Pt(18)
                p2.font.color.rgb = text_color
                p2.alignment = PP_ALIGN.CENTER
                p2.space_before = Pt(24)

        else:
            # Content / two-column
            # Title
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(1))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_title
            p.font.size = Pt(30)
            p.font.bold = True
            p.font.color.rgb = heading_color

            # Bullets
            if bullets:
                content_width = Inches(11.733) if layout != "two-column" else Inches(6)
                txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), content_width, Inches(5))
                tf2 = txBox2.text_frame
                tf2.word_wrap = True
                for i, b in enumerate(bullets):
                    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
                    p.text = f"•  {b}"
                    p.font.size = Pt(18)
                    p.font.color.rgb = text_color
                    p.space_before = Pt(8)
                    p.space_after = Pt(4)

    # Save to bytes
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


class ExportPPTXRequest(BaseModel):
    generation_id: str


@router.post("/export/pptx")
async def export_pptx(
    req: ExportPPTXRequest,
    tg_user: dict = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Export presentation as PowerPoint PPTX file."""
    tg_id = tg_user["id"]

    result = await db.execute(
        select(Generation).where(
            Generation.id == req.generation_id,
            Generation.type == "presentation",
            Generation.user_tg_id == tg_id,
        )
    )
    gen = result.scalar_one_or_none()
    if not gen:
        raise HTTPException(404, "Презентация не найдена")

    slides = json.loads(gen.result_text) if gen.result_text else []
    meta = gen.metadata_ or {}
    style = meta.get("style", "modern")
    pptx_title = gen.prompt or "Presentation"

    try:
        pptx_bytes = build_pptx(slides, style, pptx_title)
    except Exception as e:
        logger.error(f"PPTX build error: {e}")
        raise HTTPException(500, "Ошибка генерации PPTX")

    filename = re.sub(r'[^\w\s-]', '', pptx_title)[:50].strip() or "presentation"

    return StreamingResponse(
        io.BytesIO(pptx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}.pptx"'},
    )
